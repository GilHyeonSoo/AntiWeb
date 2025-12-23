"""
PDF Processor Module
Converts PDF files to images and processes them with Gemini API for OCR.
"""
import os
import io
import base64
import tempfile
from pathlib import Path

try:
    from pdf2image import convert_from_bytes
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    print("WARNING: pdf2image not installed. PDF to image conversion will not work.")

try:
    import google.generativeai as genai
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False
    print("WARNING: google-generativeai not installed. Gemini API will not work.")

from PIL import Image




def check_dependencies():
    """Check if all required dependencies are available."""
    issues = []
    if not PDF2IMAGE_AVAILABLE:
        issues.append("pdf2image not installed (pip install pdf2image)")
    if not GEMINI_AVAILABLE:
        issues.append("google-generativeai not installed (pip install google-generativeai)")
    return issues


def pdf_to_images(pdf_bytes: bytes, dpi: int = 150) -> list:
    """
    Convert PDF bytes to a list of PIL Image objects.
    
    Args:
        pdf_bytes: The PDF file as bytes
        dpi: Resolution for conversion (default 150 for balance of quality/speed)
    
    Returns:
        List of PIL Image objects, one per page
    """
    if not PDF2IMAGE_AVAILABLE:
        raise ImportError("pdf2image is not installed. Please run: pip install pdf2image")
    
    # Try common Poppler installation paths on Windows
    poppler_path = None
    if os.name == 'nt':  # Windows
        possible_paths = [
            r"C:\poppler-25.12.0\Library\bin",
            r"C:\Program Files\poppler\Library\bin",
            r"C:\Program Files\poppler-24.02.0\Library\bin",
            r"C:\poppler\Library\bin",
        ]
        for path in possible_paths:
            if os.path.exists(path):
                poppler_path = path
                break
    
    # Convert PDF to images
    if poppler_path:
        images = convert_from_bytes(pdf_bytes, dpi=dpi, poppler_path=poppler_path)
    else:
        images = convert_from_bytes(pdf_bytes, dpi=dpi)
    return images


def image_to_base64(image: Image.Image, format: str = "PNG") -> str:
    """Convert PIL Image to base64 string."""
    buffer = io.BytesIO()
    image.save(buffer, format=format)
    return base64.b64encode(buffer.getvalue()).decode('utf-8')


def extract_text_with_gemini(image: Image.Image, api_key: str) -> dict:
    """
    Extract text from an image using Gemini 1.5 Flash API.
    
    Args:
        image: PIL Image object
        api_key: Gemini API key
    
    Returns:
        Dictionary with extracted text and metadata
    """
    if not GEMINI_AVAILABLE:
        raise ImportError("google-generativeai is not installed")
    
    # Configure Gemini
    genai.configure(api_key=api_key)
    
    # Use Gemini 1.5 Flash model
    model = genai.GenerativeModel('gemini-2.0-flash-001')
    
    # Prepare the prompt
    prompt = """이 이미지의 내용을 다음 규칙에 따라 추출해주세요:

1. 모든 텍스트를 정확하게 추출합니다.
2. 수학 수식은 LaTeX 문법으로 변환합니다:
   - 인라인 수식: $수식$
   - 블록 수식: $$수식$$
3. 문제 번호, 지문, 보기를 구분하여 마크다운 형식으로 정리합니다.
4. 표가 있으면 마크다운 테이블로 변환합니다.
5. 그림/도형이 있으면 [그림: 설명] 형태로 표시합니다.

출력 형식:
- 마크다운 형식으로 깔끔하게 정리
- 원본의 구조와 순서를 유지
- 수식은 반드시 LaTeX 형식 사용"""

    # Generate content
    response = model.generate_content([prompt, image])
    
    return {
        'text': response.text,
        'success': True
    }


def process_pdf(pdf_bytes: bytes, api_key: str) -> dict:
    """
    Process an entire PDF file and extract text from all pages.
    
    Args:
        pdf_bytes: The PDF file as bytes
        api_key: Gemini API key
    
    Returns:
        Dictionary with all extracted text and metadata
    """
    # Check dependencies
    issues = check_dependencies()
    if issues:
        return {
            'success': False,
            'error': 'Missing dependencies: ' + ', '.join(issues),
            'text': ''
        }
    
    try:
        # Convert PDF to images
        images = pdf_to_images(pdf_bytes)
        
        if not images:
            return {
                'success': False,
                'error': 'PDF에서 페이지를 추출할 수 없습니다.',
                'text': ''
            }
        
        # Process each page
        all_text = []
        for i, image in enumerate(images):
            page_num = i + 1
            
            try:
                result = extract_text_with_gemini(image, api_key)
                if result['success']:
                    all_text.append(f"## 페이지 {page_num}\n\n{result['text']}")
                else:
                    all_text.append(f"## 페이지 {page_num}\n\n[오류: 텍스트 추출 실패]")
            except Exception as e:
                all_text.append(f"## 페이지 {page_num}\n\n[오류: {str(e)}]")
        
        # Combine all pages
        combined_text = '\n\n---\n\n'.join(all_text)
        
        return {
            'success': True,
            'text': combined_text,
            'page_count': len(images)
        }
        
    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'text': ''
        }


# ============ PPTX Processing with Gemini OCR ============

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("WARNING: python-pptx not installed. PPTX processing will not work.")


def pptx_to_images(pptx_bytes: bytes, dpi: int = 150) -> list:
    """
    Convert PPTX slides to images using python-pptx and PIL.
    
    Args:
        pptx_bytes: The PPTX file as bytes
        dpi: Resolution for conversion
        
    Returns:
        List of PIL Image objects, one per slide
    """
    if not PPTX_AVAILABLE:
        return []
    
    import io
    import subprocess
    import tempfile
    import os
    
    # Save PPTX to temp file
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_pptx:
        tmp_pptx.write(pptx_bytes)
        tmp_pptx_path = tmp_pptx.name
    
    images = []
    
    try:
        # Try using LibreOffice to convert PPTX to PDF first
        with tempfile.TemporaryDirectory() as tmp_dir:
            # Convert PPTX to PDF using LibreOffice (if available)
            try:
                # Try common LibreOffice paths on Windows
                soffice_paths = [
                    r'C:\Program Files\LibreOffice\program\soffice.exe',
                    r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
                ]
                
                soffice_cmd = None
                for path in soffice_paths:
                    if os.path.exists(path):
                        soffice_cmd = path
                        print(f"✅ Found LibreOffice at: {path}")
                        break
                
                if soffice_cmd:
                    print(f"🔄 Converting PPTX to PDF with LibreOffice...")
                    subprocess.run([
                        soffice_cmd, '--headless', '--convert-to', 'pdf',
                        '--outdir', tmp_dir, tmp_pptx_path
                    ], check=True, capture_output=True, timeout=120)
                    
                    # Find the generated PDF
                    pdf_path = os.path.join(tmp_dir, os.path.basename(tmp_pptx_path).replace('.pptx', '.pdf'))
                    
                    if os.path.exists(pdf_path):
                        print(f"✅ PPTX converted to PDF successfully")
                        with open(pdf_path, 'rb') as f:
                            pdf_bytes = f.read()
                        images = pdf_to_images(pdf_bytes, dpi)
                    else:
                        print(f"❌ PDF file not found after conversion")
                else:
                    print(f"❌ LibreOffice not found")
            except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired) as e:
                # LibreOffice not available or failed
                print(f"❌ LibreOffice conversion failed: {e}")
                pass
    finally:
        # Clean up temp file
        if os.path.exists(tmp_pptx_path):
            os.unlink(tmp_pptx_path)
    
    return images


def process_pptx(pptx_bytes: bytes, api_key: str) -> dict:
    """
    Process a PPTX file and extract text using Gemini OCR.
    Falls back to direct text extraction if image conversion fails.
    
    Args:
        pptx_bytes: The PPTX file as bytes
        api_key: Gemini API key
        
    Returns:
        Dictionary with success status, text content, and slide count
    """
    if not PPTX_AVAILABLE:
        return {
            'success': False,
            'error': 'python-pptx not installed',
            'text': ''
        }
    
    try:
        # First try to convert to images and use Gemini OCR
        images = pptx_to_images(pptx_bytes)
        
        if images:
            # Use Gemini OCR like PDF
            all_text = []
            for i, image in enumerate(images):
                slide_num = i + 1
                
                try:
                    result = extract_text_with_gemini(image, api_key)
                    if result['success']:
                        all_text.append(f"## 슬라이드 {slide_num}\n\n{result['text']}")
                    else:
                        all_text.append(f"## 슬라이드 {slide_num}\n\n[오류: 텍스트 추출 실패]")
                except Exception as e:
                    all_text.append(f"## 슬라이드 {slide_num}\n\n[오류: {str(e)}]")
            
            combined_text = '\n\n---\n\n'.join(all_text)
            
            return {
                'success': True,
                'text': combined_text,
                'slide_count': len(images)
            }
        else:
            # Fallback: Direct text extraction (no OCR)
            import io
            pptx_file = io.BytesIO(pptx_bytes)
            prs = Presentation(pptx_file)
            
            all_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"## 슬라이드 {slide_num}"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:
                    all_text.append('\n\n'.join(slide_text))
            
            combined_text = '\n\n---\n\n'.join(all_text)
            
            return {
                'success': True,
                'text': combined_text,
                'slide_count': len(prs.slides),
                'fallback': True  # Indicates direct extraction was used
            }
        
    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'text': ''
        }


# ============ DOCX Processing ============

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("WARNING: python-docx not installed. DOCX processing will not work.")


def extract_docx_text(docx_bytes: bytes) -> dict:
    """
    Extract text from a DOCX file.
    
    Args:
        docx_bytes: The DOCX file as bytes
        
    Returns:
        Dictionary with success status and text content
    """
    if not DOCX_AVAILABLE:
        return {
            'success': False,
            'error': 'python-docx not installed',
            'text': ''
        }
    
    try:
        import io
        docx_file = io.BytesIO(docx_bytes)
        doc = Document(docx_file)
        
        all_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                all_text.append(para.text.strip())
        
        combined_text = '\n\n'.join(all_text)
        
        return {
            'success': True,
            'text': combined_text,
            'paragraph_count': len(all_text)
        }
        
    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'text': ''
        }
